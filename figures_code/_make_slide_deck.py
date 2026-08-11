"""Build the WRC SARA 2026 oral-presentation deck for paper WRC26_0024.

Built on the conference's own template (paper/WRC_oral_template.pptx, downloaded
from the SARA site) so the title slide, the blue header bar, the WRC logo and the
slide numbering are exactly what the organisers specified. The template's four
sample slides are dropped; the closing "Thank You" slide is kept and moved to the
end.

Every slide carries its narration in the speaker-notes pane, so the same file
drives the live talk and the recorded video.

Usage:
    python -m paper._make_slide_deck
Output:
    paper/FlyAda_WRC_SARA_2026_oral.pptx
"""
from __future__ import annotations

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from paper._slide_script import SCRIPT, mmss, timings   # noqa: E402

# Derived once: slide-key -> "[m:ss – m:ss]" cue prepended to the speaker notes.
# The Thank You slide is off the clock, so it gets a label instead of a range.
_CUES = {k: ("[after the talk]" if k == "thanks" else f"[{mmss(a)} – {mmss(b)}]")
         for k, _n, _t, _w, a, b, _backup in timings()}

TEMPLATE = ROOT / "paper" / "WRC_oral_template.pptx"
FIGS = ROOT / "paper" / "figures"
SLIDES = FIGS / "slides"
OUT = ROOT / "paper" / "FlyAda_WRC_SARA_2026_oral.pptx"

SW, SH = Inches(13.3333), Inches(7.5)
HEADER_H = Inches(0.6)

FONT = "Segoe UI"
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
C_VAN = RGBColor(0x2A, 0x78, 0xD6)
C_FLY = RGBColor(0xEB, 0x68, 0x34)
BAND = RGBColor(0xF3, 0xF6, 0xFB)

CONTENT_LAYOUT = 1          # 标题和内容 — carries the blue bar, logo, slide number


# ------------------------------------------------------------------ helpers ---
def drop_slides(prs, drop_idx):
    """Delete template slides by index.

    Call this only AFTER every new slide has been added: python-pptx assigns
    slideN.xml part names from the set of currently reachable slides, so
    deleting first makes the freed numbers collide with the survivors and the
    saved package will not open.
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for i in sorted(drop_idx, reverse=True):
        rId = slides[i].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        xml_slides.remove(slides[i])


def move_slide(prs, src, dst):
    """Move the slide at index `src` to index `dst` (indices before the move)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    node = slides[src]
    xml_slides.remove(node)
    remaining = list(xml_slides)
    if dst >= len(remaining):
        xml_slides.append(node)
    else:
        remaining[dst].addprevious(node)


def set_text(tf, text, size, *, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font=FONT, space_after=6, line=None):
    tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tf


def add_textbox(slide, x, y, w, h, text, size, **kw):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.margin_left = tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = tb.text_frame.margin_bottom = Emu(0)
    set_text(tb.text_frame, text, size, **kw)
    return tb


def add_bullets(slide, x, y, w, h, items, size=17, color=INK2, gap=11):
    """Bulleted body copy — a dot glyph plus a hanging indent, drawn by hand so
    the template's own list styling doesn't fight the layout."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.18
        bold_head = None
        body = item
        if isinstance(item, tuple):
            bold_head, body = item
        r = p.add_run()
        r.text = "▪  "
        r.font.size = Pt(size)
        r.font.color.rgb = C_FLY
        r.font.name = FONT
        if bold_head:
            rb = p.add_run()
            rb.text = bold_head
            rb.font.size = Pt(size)
            rb.font.bold = True
            rb.font.color.rgb = INK
            rb.font.name = FONT
        r2 = p.add_run()
        r2.text = body
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = FONT
    tb.name = "BULLETS"
    return tb


def add_title(slide, text):
    """Fill the layout's title placeholder — that is the blue header bar."""
    ph = slide.shapes.title
    ph.left, ph.top, ph.width, ph.height = Emu(0), Emu(0), SW, HEADER_H
    tf = ph.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.22)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(24)
    r.font.bold = False
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = FONT
    return ph


def add_takeaway(slide, text, y=Inches(6.36), color=C_FLY):
    """The one-sentence 'so what' strip along the bottom of a slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.42), y, Inches(12.5), Inches(0.78))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BAND
    bar.line.color.rgb = color
    bar.line.width = Pt(1.5)
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = INK
    r.font.name = FONT
    bar.name = "TAKEAWAY"
    return bar


def add_picture_fit(slide, img, x, y, w, h):
    """Place an image centred inside the (x, y, w, h) box, preserving aspect."""
    from PIL import Image
    iw, ih = Image.open(img).size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    pic = slide.shapes.add_picture(str(img), int(x + (w - nw) / 2),
                                   int(y + (h - nh) / 2), nw, nh)
    pic.name = "FIGURE"
    return pic


def add_movie_fit(slide, mp4, poster, x, y, w, h):
    from PIL import Image
    iw, ih = Image.open(poster).size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    return slide.shapes.add_movie(str(mp4), int(x + (w - nw) / 2),
                                  int(y + (h - nh) / 2), nw, nh,
                                  poster_frame_image=str(poster),
                                  mime_type="video/mp4")


def new_slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
    # The layout ships a body placeholder we never use — drop it.
    for shp in list(s.shapes):
        if shp.is_placeholder and shp.placeholder_format.idx == 1:
            shp._element.getparent().remove(shp._element)
    add_title(s, title)
    return s


def notes(slide, key):
    slide.notes_slide.notes_text_frame.text = (
        f"{_CUES[key]}\n{SCRIPT[key].strip()}\n")


# ================================================================== the deck ===
def build():
    prs = Presentation(str(TEMPLATE))
    prs.slide_width, prs.slide_height = SW, SH

    # ---------------------------------------------------------- title slide ---
    t = prs.slides[0]
    texts = {
        "This Is The Oral Paper Title with Long Long Long Long Title":
            "FlyAda: Belief-State Adaptation for Diffusion Policies\nunder Partial Observation",
        "Author1, Author2, Author3, Author4, Author5, Author6, Author7, Author8, "
        "Author9, Author10, Author11, Author12, Author13": "Lingxue Lyu",
        "Xxx xxx xxx Lab., xxx University":
            "University of Pennsylvania, Philadelphia, PA, United States",
    }
    for shp in t.shapes:
        if not shp.has_text_frame:
            continue
        cur = shp.text_frame.text.strip()
        for k, v in texts.items():
            if cur == k.strip():
                tf = shp.text_frame
                # Keep the template's paragraph formatting; swap the words only.
                first = tf.paragraphs[0]
                keep = first.runs[0] if first.runs else None
                for p in list(tf.paragraphs)[1:]:
                    p._p.getparent().remove(p._p)
                for r in list(first.runs)[1:]:
                    r._r.getparent().remove(r._r)
                if keep is not None:
                    keep.text = v.replace("\n", "  ")
    notes(prs.slides[0], "title")

    # ===================================================== MAIN LINE (2-10) ===
    # Ten slides: one question, one failure, one method, three results. Anything
    # that is evidence for a claim the audience already accepted lives in backup.

    # -------------------------------------------------------- 2 the question --
    s = new_slide(prs, "1 · The question")
    add_bullets(s, Inches(0.7), Inches(1.30), Inches(7.5), Inches(3.6), [
        ("Action-chunk diffusion policies ", "are a strong imitation-learning "
         "recipe for continuous control."),
        ("The common assumption: ", "they should execute poorly when test-time "
         "dynamics shift away from training, because the denoiser conditions on "
         "a single observation."),
        ("Our prior work ", "found that a small online-updated latent closed "
         "exactly that kind of mismatch gap — and we set out to port that recipe "
         "to the diffusion setting."),
    ], size=18, gap=16)
    add_textbox(s, Inches(0.7), Inches(5.05), Inches(7.5), Inches(1.5),
                "When does a small online-updated\nlatent actually help?",
                26, bold=True, color=C_FLY, line=1.25)
    add_picture_fit(s, SLIDES / "fig_story_arc.png",
                    Inches(8.45), Inches(0.80), Inches(4.65), Inches(6.0))
    notes(s, "question")

    # ------------------------------------------------------------- 3 the task --
    s = new_slide(prs, "2 · The task")
    add_movie_fit(s, SLIDES / "chunk_replan.mp4",
                  SLIDES / "chunk_replan_still.png",
                  Inches(0.25), Inches(0.95), Inches(7.35), Inches(4.45))
    add_picture_fit(s, SLIDES / "fig_task_card.png",
                    Inches(7.85), Inches(0.75), Inches(5.3), Inches(5.9))
    add_textbox(s, Inches(0.5), Inches(5.65), Inches(7.0), Inches(0.5),
                "UAV goal-reaching · 50 Hz · plan 8 actions, execute 4, re-plan",
                14, color=INK2, align=PP_ALIGN.CENTER)
    notes(s, "setup")

    # ------------------------------------------- 4 full obs is already robust --
    s = new_slide(prs, "3 · Under full observation, vanilla is already robust")
    add_picture_fit(s, SLIDES / "fig_fullobs.png",
                    Inches(0.5), Inches(0.95), Inches(12.3), Inches(5.2))
    add_takeaway(s, "So the failure is not caused by dynamics shift itself — "
                    "there is nothing here for an adaptation latent to do.")
    notes(s, "fullobs")

    # --------------------------------------------------- 5 hiding velocity ----
    s = new_slide(prs, "4 · Hide velocity, and it breaks")
    add_movie_fit(s, SLIDES / "overshoot.mp4", SLIDES / "overshoot_still.png",
                  Inches(0.35), Inches(0.75), Inches(9.5), Inches(5.5))
    add_bullets(s, Inches(10.05), Inches(1.5), Inches(3.0), Inches(4.4), [
        ("Same seed, same backbone. ", "Only the observation differs."),
        ("It accelerates correctly — ", "it just cannot tell when to brake."),
        ("0.3% ", "success across the 12-condition sweep."),
    ], size=16, gap=22)
    add_takeaway(s, "Closer in position (0.40 m vs 0.50 m) and still a failure, "
                    "because it arrives at 1.45 m/s through a 1 m/s stop tolerance.")
    notes(s, "partial")

    # ------------------------------------------------ 6 why frame-stack fails --
    s = new_slide(prs, "5 · Why more history is not the answer")
    add_textbox(s, Inches(0.7), Inches(0.95), Inches(12.1), Inches(1.4),
                "The history contains the velocity.  "
                "The imitation loss never asks the model to extract it.",
                26, bold=True, color=C_FLY, line=1.3)
    add_picture_fit(s, SLIDES / "fig_ambiguity.png",
                    Inches(0.3), Inches(2.35), Inches(7.9), Inches(3.9))
    add_picture_fit(s, FIGS / "flyada_partial_loss.png",
                    Inches(8.55), Inches(2.20), Inches(4.4), Inches(4.05))
    add_textbox(s, Inches(8.55), Inches(6.30), Inches(4.4), Inches(0.5),
                "All three training losses track each other.",
                13, color=INK2, align=PP_ALIGN.CENTER)
    add_takeaway(s, "A 3-frame stack carries the same information frame by frame "
                    "— and still reaches only 1.1%.", y=Inches(6.60))
    notes(s, "why")

    # ------------------------------------------------------------- 7 FlyAda ---
    s = new_slide(prs, "6 · FlyAda — a small observer with an auxiliary loss")
    add_picture_fit(s, SLIDES / "fig_method.png",
                    Inches(0.35), Inches(0.72), Inches(12.65), Inches(5.55))
    add_takeaway(s, "Same observation, same denoiser, same demonstrations. "
                    "One 73 K-parameter head, one EMA update, one auxiliary loss.")
    notes(s, "method")

    # -------------------------------------------------------- 8 main result ---
    s = new_slide(prs, "7 · Main result — velocity hidden")
    add_picture_fit(s, SLIDES / "fig_sweep_grid.png",
                    Inches(0.3), Inches(0.72), Inches(6.5), Inches(5.5))
    add_movie_fit(s, SLIDES / "taskB_race_nominal.mp4",
                  SLIDES / "taskB_race_nominal_still.png",
                  Inches(6.85), Inches(1.05), Inches(6.3), Inches(4.6))
    add_textbox(s, Inches(6.85), Inches(5.75), Inches(6.3), Inches(0.5),
                "Task B — 3 waypoints, 800 steps, no retraining",
                13.5, color=INK2, align=PP_ALIGN.CENTER)
    add_takeaway(s, "0.3% → 100% on the sweep — and 1.00 / 0.93 on a 3-waypoint "
                    "chain the policy was never trained for.")
    notes(s, "results")

    # ------------------------------------------------------- 9 sim-to-sim -----
    s = new_slide(prs, "8 · It transfers to a 6-DoF body")
    add_movie_fit(s, SLIDES / "mujoco_flight.mp4",
                  SLIDES / "mujoco_flight_still.png",
                  Inches(0.25), Inches(0.70), Inches(8.6), Inches(5.5))
    add_picture_fit(s, SLIDES / "fig_mujoco_bars.png",
                    Inches(9.0), Inches(0.95), Inches(4.1), Inches(3.6))
    add_textbox(s, Inches(9.0), Inches(4.80), Inches(4.15), Inches(1.4),
                "1 kg · 0.15 m arms · four rotors\n"
                "500 Hz attitude-rate PID under a 50 Hz policy\n"
                "No retuning, no retraining — same checkpoints.",
                12.5, color=INK2, line=1.3, space_after=4)
    add_takeaway(s, "The only policy with non-trivial transfer: 0.40 nominal / "
                    "0.85 hard, against 0.00 for both baselines.")
    notes(s, "mujoco")

    # --------------------------------------------------- 10 what does the work -
    s = new_slide(prs, "9 · What is doing the work")
    add_picture_fit(s, SLIDES / "fig_ablation_compact.png",
                    Inches(0.25), Inches(0.80), Inches(6.3), Inches(5.4))
    add_movie_fit(s, SLIDES / "belief_tracking.mp4",
                  SLIDES / "belief_tracking_still.png",
                  Inches(6.75), Inches(0.78), Inches(6.4), Inches(3.05))
    add_picture_fit(s, SLIDES / "fig_latent_pca.png",
                    Inches(6.85), Inches(3.90), Inches(3.6), Inches(2.35))
    add_bullets(s, Inches(10.6), Inches(4.15), Inches(2.55), Inches(2.1), [
        ("Linear probe ", "z → v :  R² = 0.992"),
        ("Classifier ", "z → condition :  24%  (chance 20%)"),
    ], size=12.5, gap=10)
    add_takeaway(s, "Continuous online updating is the load-bearing part — and "
                    "the latent is the missing state channel, not a regime label.")
    notes(s, "mechanism")

    # ------------------------------------------------ 10 caveats + conclusion --
    s = new_slide(prs, "10 · Caveats and conclusion")
    add_bullets(s, Inches(0.75), Inches(0.95), Inches(11.9), Inches(1.4), [
        ("Simulation only. ", "No hardware."),
        ("The auxiliary loss needs a true-velocity target at training time — ",
         "mild for velocity, but it does not obviously extend to quantities you "
         "cannot instrument."),
    ], size=16.5, gap=8)

    add_textbox(s, Inches(0.9), Inches(2.72), Inches(11.6), Inches(2.5),
                "Diffusion policies are robust when they can see the state.\n"
                "They fail when a state channel is hidden — and more history "
                "does not fix it.\n"
                "FlyAda fixes it by making a small latent explicitly learn the "
                "missing channel.",
                23, bold=True, color=INK, line=1.35, space_after=10)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.42), Inches(2.50), Inches(12.5), Inches(2.42))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BAND
    bar.line.color.rgb = C_FLY
    bar.line.width = Pt(1.5)
    bar.shadow.inherit = False
    bar.name = "PANEL"
    s.shapes._spTree.remove(bar._element)
    s.shapes._spTree.insert(2, bar._element)      # behind the text

    add_picture_fit(s, SLIDES / "fig_scoreboard.png",
                    Inches(0.45), Inches(5.10), Inches(12.45), Inches(1.85))
    add_textbox(s, Inches(0.75), Inches(7.02), Inches(11.9), Inches(0.4),
                "Code, configs and trained-policy artefacts: "
                "github.com/…/flyada     ·     lingxuelyu@alumni.upenn.edu",
                11.5, color=INK2, align=PP_ALIGN.CENTER)
    notes(s, "conclusion")

    # ---------------------------------------------------------- 11 thank you --
    notes(prs.slides[5], "thanks")     # the template's closing slide

    # ---------------------------- drop the template's samples, order the deck --
    # After dropping the four samples the template's Thank You sits at index 1,
    # ahead of everything we added; it belongs last.
    drop_slides(prs, {1, 2, 3, 4})
    move_slide(prs, 1, 11)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"saved -> {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
